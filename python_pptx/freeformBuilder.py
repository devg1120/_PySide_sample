
import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches


import subprocess

prs = pptx.Presentation()


#########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


slide = prs.slides.add_slide(prs.slide_layouts[6]) 

"""
shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "いろいろ"
run.font.color.rgb = RGBColor(0, 0, 255)


# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#slideshapes-objects
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype

a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(3), height=Cm(3) ,left=Cm(3), top=Cm(3))
l = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE,       width=Cm(5), height=Cm(5) ,left=Cm(5), top=Cm(5))

# https://github.com/scanny/python-pptx/blob/master/docs/dev/analysis/shp-freeform.rst

freeform_builder = slide.shapes.build_freeform(Cm(10), Cm(10))
freeform_builder.add_line_segments((
     (Cm(10),   Cm(10)),
     (Cm(9),   Cm(10)),
     (Cm(9.5), Cm(9)),
 ),
 close = False
 )
freeform_shape = freeform_builder.convert_to_shape()
"""
"""
a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(3), top=Cm(12))
b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(11), top=Cm(10))

pg = a.text_frame.paragraphs[0]
pg.text ="A"
pg.alignment = PP_ALIGN.CENTER

pg = b.text_frame.paragraphs[0]
pg.text ="B"
pg.alignment = PP_ALIGN.CENTER

#line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(2), Cm(2), Cm(2), Cm(2))
line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(a, 3)
line1.end_connect(b, 1)
"""

c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(14), top=Cm(12))
d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(22), top=Cm(10))

pg = c.text_frame.paragraphs[0]
pg.text ="C"
pg.alignment = PP_ALIGN.CENTER

pg = d.text_frame.paragraphs[0]
pg.text ="D"
pg.alignment = PP_ALIGN.CENTER

line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(c, 3)
line1.end_connect(d, 2)
line1.element.spPr.prstGeom.rewrite_guides([('adj1', -635)])

e = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(14), top=Cm(15))
f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(22), top=Cm(17))

pg = e.text_frame.paragraphs[0]
pg.text ="E"
pg.alignment = PP_ALIGN.CENTER

pg = f.text_frame.paragraphs[0]
pg.text ="F"
pg.alignment = PP_ALIGN.CENTER


line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2) )
#line1.element.flipH = True
#line1.element.flipV = True
line1.element.rot = 16200000
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(e, 0)
line1.end_connect(f, 0)

#https://github.com/scanny/python-pptx/blob/master/src/pptx/oxml/shapes/shared.py
#line1.element.flipH = True
#line1.element.flipV = True
#line1.element.rot = 16200000


#e.element.spPr.prstGeom.rewrite_guides([("adj1", 25000), ("adj2", 25000), ("adj3", 25000), ("adj4", 43750)])
#print(e.element.spPr.prstGeom.gd_lst)
#e.element.flipH = True
#e.element.rot = 1620
print(type(line1.element))
print(type(line1.element.nvCxnSpPr))
print(type(line1.element.nvCxnSpPr.cNvPr)) #<class 'pptx.oxml.shapes.shared.CT_NonVisualDrawingProps'>

print(type(line1.element.nvCxnSpPr.cNvCxnSpPr))  #<class 'pptx.oxml.shapes.connector.CT_NonVisualConnectorProperties'>
print(type(line1.element.nvCxnSpPr.cNvCxnSpPr.stCxn))  #<class 'pptx.oxml.shapes.connector.CT_Connection'>
print(type(line1.element.nvCxnSpPr.cNvCxnSpPr.endCxn)) #<class 'pptx.oxml.shapes.connector.CT_Connection'>

print(type(line1.element.nvCxnSpPr.nvPr))
print(type(line1.element.spPr))
line1.element.spPr.prstGeom.rewrite_guides([('adj1', 0)])


#line1.element.spPr.prstGeom.rewrite_guides([("adj1", 50000), ("adj2", 50000)])
#line1.element.spPr.prstGeom.rewrite_guides([('adj2', 3500)])

##########################################################################################################
##########################################################################################################
##########################################################################################################


prs.save('ffb1.pptx')


subprocess.Popen(['start', 'ffb1.pptx'], shell=True)



