import os
import signal
import psutil

import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.util import Cm,Pt,Inches

from pptx.util import lazyproperty


default_fill_color = RGBColor(100, 149, 237)
default_line_color = RGBColor( 65, 105, 225)
default_line_weight = 1
default_text_color = RGBColor(0,0,0)
default_text_size = 16
default_text_align =  PP_ALIGN.CENTER
default_text_bold = True
default_text_italic = False


class Point():
   def  __init__(self,  x = 0 , y = 0 ):
       self.x = x
       self.y = y


class Shape():
   def  __init__(self, slide , shape_type, 
           width  , height , left , top , 
           fill_color, line_color, line_weight ,
           text, text_color, text_size, text_align , text_bold, text_italic):

         self.slide  = slide
         self.shape_type = shape_type
         self.width  = width
         self.height = height
         self.left   = left
         self.top    = top

         if not slide == None:
            #self.shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            self.shape = slide.shapes.add_shape(shape_type, 
                 width=Pt(width), height=Pt(height) ,left=Pt(left), top=Pt(top))
            self.shape.fill.solid()
            #self.shape.fill.fore_color.rgb = RGBColor(250, 100, 100)
            self.shape.fill.fore_color.rgb = fill_color
            self.shape.line.fill.background()
            self.shape.line.fill.solid()
            self.shape.line.fill.fore_color.rgb = line_color
            self.shape.line.width = Pt(line_weight)
            #self.shape.text =  text
            pg = self.shape.text_frame.paragraphs[0]
            pg.text = text

            #pg.alignment = PP_ALIGN.CENTER
            #pg.font.color.rgb =  RGBColor(0,0,0)
            #pg.font.size =  Pt(10)
            #pg.font.bold =  True
            #pg.font.italic =  True
            pg.alignment =  text_align
            pg.font.color.rgb = text_color
            pg.font.size =  Pt(text_size)
            pg.font.bold =  text_bold
            pg.font.italic =  text_italic

class Rectangle(Shape):
   def  __init__(self, slide , 
           width = 0 , height = 0 , left = 0, top = 0,
           text = "",
           fill_color  = default_fill_color,
           line_color  = default_line_color,
           line_weight = default_line_weight,
           text_color  = default_text_color,
           text_size   = default_text_size,
           text_align  = default_text_align,
           text_bold  = default_text_bold,
           text_italic  = default_text_italic,
           ):
      super().__init__(slide , MSO_SHAPE.RECTANGLE, 
                       width , height , left , top ,
                       fill_color, line_color, line_weight,
                       text, text_color, text_size, text_align,text_bold, text_italic)

class RoundedRectangle(Shape):
   def  __init__(self, slide , 
           width = 0 , height = 0 , left = 0, top = 0,
           text = "",
           fill_color  = default_fill_color,
           line_color  = default_line_color,
           line_weight = default_line_weight,
           text_color  = default_text_color,
           text_size   = default_text_size,
           text_align  = default_text_align,
           text_bold  = default_text_bold,
           text_italic  = default_text_italic,
           rounded = None
           ):
      super().__init__(slide , MSO_SHAPE.ROUNDED_RECTANGLE, 
                       #width , height , left , top , text , fill_color, line_color, line_weight)
                       width , height , left , top ,
                       fill_color, line_color, line_weight,
                       text, text_color, text_size, text_align,text_bold, text_italic)
      if not rounded == None:
         self.shape.adjustments[0] = rounded

class Cube(Shape):
   def  __init__(self, slide , 
           width = 0 , height = 0 , left = 0, top = 0,
           text = "",
           fill_color = default_fill_color,
           line_color = default_line_color,
           line_weight = default_line_weight,
           text_color  = default_text_color,
           text_size   = default_text_size,
           text_align  = default_text_align,
           text_bold  = default_text_bold,
           text_italic  = default_text_italic,
           depth = None
           ):
      super().__init__(slide , MSO_SHAPE.CUBE, 
                       #width , height , left , top , text , fill_color, line_color, line_weight)
                       width , height , left , top ,
                       fill_color, line_color, line_weight,
                       text, text_color, text_size, text_align,text_bold, text_italic)
      if not depth == None:
         self.shape.adjustments[0] = depth

class Arc(Shape):
   def  __init__(self, slide , 
           width = 0 , height = 0 , left = 0, top = 0,
           text = "",
           fill_color = default_fill_color,
           line_color = default_line_color,
           line_weight = default_line_weight,
           text_color  = default_text_color,
           text_size   = default_text_size,
           text_align  = default_text_align,
           text_bold  = default_text_bold,
           text_italic  = default_text_italic,
           radius0 = None,
           radius1 = None,
           ):
      super().__init__(slide , MSO_SHAPE.ARC, 
                       #width , height , left , top , text , fill_color, line_color, line_weight)
                       width , height , left , top ,
                       fill_color, line_color, line_weight,
                       text, text_color, text_size, text_align,text_bold, text_italic)
      print("adj len",  len(self.shape.adjustments))
      print("adj[0]",   self.shape.adjustments[0])
      print("adj[1]",   self.shape.adjustments[1])

      if not radius0 == None:
         self.shape.adjustments[0] = radius0
      if not radius1 == None:
         self.shape.adjustments[1] = radius1

class Oval(Shape):
   def  __init__(self, slide , 
           width = 0 , height = 0 , left = 0, top = 0,
           text = "",
           fill_color = default_fill_color,
           line_color = default_line_color,
           line_weight = default_line_weight,
           text_color  = default_text_color,
           text_size   = default_text_size,
           text_align  = default_text_align,
           text_bold  = default_text_bold,
           text_italic  = default_text_italic,
           ):
      super().__init__(slide , MSO_SHAPE.OVAL, 
                       #width , height , left , top , text , fill_color, line_color, line_weight)
                       width , height , left , top ,
                       fill_color, line_color, line_weight,
                       text, text_color, text_size, text_align,text_bold, text_italic)

class TextBox():
   def  __init__(self, slide = None , width = 0 , height = 0 , left = 0, top = 0, text = "NoSet"):
         self.slide  = slide
         self.width  = width
         self.height = height
         self.left   = left
         self.top    = top
         self.text = text

         if not slide == None:
            self.shape = slide.shapes.add_textbox( 
                 width=Pt(width), height=Pt(height) ,left=Pt(left), top=Pt(top))
            self.shape.fill.solid()
            self.shape.fill.fore_color.rgb = RGBColor(250, 100, 100)
            self.shape.text = self.text

class Line():
   def  __init__(self, slide = None , start = None, end = None, pointlist = None):
         self.slide  = slide
         self.start  = start
         self.end    = end
         self.shape = None
         self.shapes = []

         if slide == None :
             retutn
         if pointlist == None:
            self.shape = slide.shapes.add_connector(
                     #MSO_CONNECTOR.STRAIGHT, 
                     #MSO_CONNECTOR.ELBOW, 
                     MSO_CONNECTOR.CURVE, 
                    Pt(start.x), 
                    Pt(start.y),
                    Pt(end.x ), 
                    Pt(end.y ),
                 )
            self.shape.line.fill.background()
            self.shape.line.fill.solid()
            self.shape.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
            self.shape.line.width = Pt(5)
            self.shape.text = "LINE1"
         else:
             for i, point in enumerate(pointlist[:-1]):
                 next_point = pointlist[i + 1]
                 shape = slide.shapes.add_connector(
                     MSO_CONNECTOR.STRAIGHT, 
                     Pt(point.x), #begin_x, 
                     Pt(point.y), #begin_y, 
                     Pt(next_point.x), #end_x, 
                     Pt(next_point.y)  #end_y
                   )
                 self.shapes.append(shape)
             self.group = slide.shapes.add_group_shape(self.shapes)


class FreeForm():
    def __int__(self, slide, pointlist = [], text = ""):
        if len(pointlist) < 2:
            return
        x = pointlist[0].x
        y = pointlist[0].y
        freeform_builder = slide.shapes.build_freeform(Pt(x),Pt(y))
        for i, point in enumerate(pointlist):
             freeform_builder.add_line_segments((
                     Pt(point.x), 
                     Pt(point.y)  
                   ), clise = False)
        self.freeform_shape = freeform_builder.convert_to_shape()
        self.freeform_shape.text = "free TEXT1"


class Table():
    def  __init__(self, slide, row, col, x, y, width, height):
        self.table = slide.shapes.add_table(row, col, Pt(x), Pt(y), Pt(width), Pt(height) ).table

    def get_table():
        return self.table

    def cell(self, row_idx, col_idx):
       return self.table.cell(row_idx, col_idx)

    def first_col(self):
        return self.table.first_col

    def first_col(self, value):
        self.table.first_col = value

    def first_row(self):
        return self.table.first_row

    def first_row(self, value):
        self.table.first_row = value

    @lazyproperty
    def columns(self):
        return self.table.columns

    @lazyproperty
    def rows(self):
        return self.table.rows
    

