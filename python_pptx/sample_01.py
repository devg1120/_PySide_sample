from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE    # 図形の定義がされているクラス
from pptx.dml.color import RGBColor       # 色を管理するクラス
from pptx.util import Cm,Pt               # 単位指定をするクラス(センチメートル, ポイント単位)

prs = Presentation()
sld0 = prs.slides.add_slide(prs.slide_layouts[6]) # 空白のスライドを追加


# [A]図形の追加と形状の変形
#-----------------------------------------------------------------------------------------------
rect0 = sld0.shapes.add_shape(		# shapeオブジェクト➀を追加
    MSO_SHAPE.ROUNDED_RECTANGLE,   	# 図形の種類を[丸角四角形]に指定
    Cm(2), Cm(2),                  	# 挿入位置の指定　左上の座標の指定
    Cm(5), Cm(3))                  	# 挿入図形の幅と高さの指定

for i in range(3):

    rect1 = sld0.shapes.add_shape(     	# shapeオブジェクト➁を追加
        MSO_SHAPE.DECAGON,              # 図形の種類を[六角形]に指定
        Cm(2+i*5), Cm(7),
        Cm(3), Cm(3))
    
    rect1.adjustments[0] = 0.5*(i+1)   # 図形の形状を変形(高さ方向)


for i in range(3):
    rect2 = sld0.shapes.add_shape(     	# shapeオブジェクト➂を追加
        MSO_SHAPE.STAR_10_POINT,        # 図形の種類を[星(10)]に指定
        Cm(2+i*5), Cm(12),
        Cm(3), Cm(3))
    
    rect2.adjustments[0] = 0.21*(i+1)   # 図形の形状を変形
    

# [B]図形の塗り潰しとテキストの設定
#-----------------------------------------------------------------------------------------------
rect0.fill.solid()                                   # shapeオブジェクト➀を単色で塗り潰す
rect0.fill.fore_color.rgb = RGBColor(250, 100, 100)  # RGB指定で色を指定

pg = rect0.text_frame.paragraphs[0]   	# shapeオブジェクト➀のTextFrameの取得
pg.text = 'ROUNDED_RECTANGLE'         	# TextFrameにテキストを設定
pg.font.size = Pt(10)                	# テキストの文字サイズを10ポイントとする


prs.save('Blog_図形の挿入.pptx')
