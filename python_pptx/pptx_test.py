
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

#プレゼンテーションを開く
prs = Presentation()

#1ページ目（「タイトル スライド」のレイアウトを指定）
slide_layout_0 = prs.slide_layouts[0]
slide_1 = prs.slides.add_slide(slide_layout_0)

#プレースホルダーの指定
s1_title = slide_1.placeholders[0]
s1_subtitle = slide_1.placeholders[1]

#textをそれぞれのプレースフォルダに設定
s1_title.text = "PythonでPowerpointを作成しよう"
s1_subtitle.text = "TeamXeppet Channel"

#2ページ目（「タイトルとコンテンツ」のレイアウトを指定）
slide_layout_1 = prs.slide_layouts[1]
slide_2 = prs.slides.add_slide(slide_layout_1)

s2_title = slide_2.placeholders[0]
s2_content = slide_2.placeholders[1]

#タイトルを緑色で追加
s2_title.text = "PythonでPowerpointを扱うには？"
pg = s2_title.text_frame.paragraphs[0]
run1 = pg.runs[0]
run1.font.color.rgb = RGBColor(0, 255, 0)

#コンテンツにテキストを追加
s2_content.text = "pip install python-pptx"+"\n"
s2_content.text += "をまずは実行しよう"

#画像の追加
#img_path = "python-pptx_install.png"
img_path = "sample_picture.jpg"
pic = slide_2.shapes.add_picture(img_path,Inches(1),Inches(3.5))

#3ページ目（「タイトルのみ」のレイアウトを指定）
slide_layout_5 = prs.slide_layouts[5]
slide_3 = prs.slides.add_slide(slide_layout_5)

s3_title = slide_3.placeholders[0]
s3_title.text = "テーブルの作成"

#テーブルに入力するデータ
col_names = ["No.","qualification","type"]
numbers = [1, 2, 3]
qnames = ["ORACLE MASTER", "Fundamental Information Technology Engineer Examination", "Linux Professional Institute Certification"]
qtypes = [ "Vendor", "Nation", "Private"]

q_data = [numbers, qnames, qtypes]

#テーブルの作成位置を指定
x, y, cx, cy = Inches(1.8), Inches(3), Inches(4), Inches(1.5)
table = slide_3.shapes.add_table(4, 3, x, y, cx, cy).table

#テーブルの列幅を指定
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(4)
table.columns[2].width = Inches(2)

#テーブルのセルへ値を挿入
for i in range(len(col_names)):
    cell = table.cell(0, i)
    cell.text = col_names[i]

for i in range(len(q_data)):
    for j in range(len(q_data[i])):
        cell = table.cell(j+1, i)
        cell.text = str(q_data[i][j])

prs.save("create_powerpnt.pptx")

