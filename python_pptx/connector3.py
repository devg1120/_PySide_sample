
import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches


import subprocess

prs = pptx.Presentation()


##########################################################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

s1 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(18), top=Cm(2), height=pptx.util.Inches(0.6))
s2 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(21), top=Cm(5), height=pptx.util.Inches(0.6))

line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line.begin_connect(s1, 3)
line.end_connect(s2, 0)



a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(3), top=Cm(1))
b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(13), top=Cm(5))

line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(2), Cm(2), Cm(2), Cm(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(a, 3)
line1.end_connect(b, 0)

c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(3), top=Cm(4))
d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(13), top=Cm(8))


line2 = slide.shapes.add_connector(MSO_CONNECTOR.CURVE, Cm(2), Cm(2), Cm(2), Cm(2))
line2.line.fill.background()
line2.line.fill.solid()
line2.line.fill.fore_color.rgb = RGBColor(0, 0, 255)
line2.begin_connect(c, 0)
line2.end_connect(d, 1)


e = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(3), top=Cm(11))
f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(13), top=Cm(13))

line2 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line2.line.fill.background()
line2.line.fill.solid()
line2.line.fill.fore_color.rgb = RGBColor(204, 0, 204)
line2.begin_connect(e, 3)
line2.end_connect(f, 0)


##########################################################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(3), top=Cm(1))
b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(13), top=Cm(5))

line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(2), Cm(2), Cm(2), Cm(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(a, 3)
line1.end_connect(b, 0)


##########################################################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6]) 


a1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(0), top=Cm(2))
b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(2), top=Cm(4))
c1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(4), top=Cm(6))
g1 = slide.shapes.add_group_shape([a1, b1, c1] )


a2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(7), top=Cm(9))
b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(9), top=Cm(11))
c2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(1), height=Cm(1) ,left=Cm(11), top=Cm(13))
g2 = slide.shapes.add_group_shape([a2, b2, c2] )


##########################################################################################################
# https://hogelog.com/python/python-powerpoint-2.html

slide = prs.slides.add_slide(prs.slide_layouts[6]) 

shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

#テキストの一部の色を赤色に変える
pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "テキスト"
run.font.color.rgb = RGBColor(255, 0, 0)

#テキストの一部を青色に変える
run = pg.add_run()
run.text = " 表示いろいろ"
run.font.color.rgb = RGBColor(0, 0, 255)


shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(3), Cm(5), Cm(5))

#shapeオブジェクトでtextを挿入
shape.text ="Python"
text_frame = shape.text_frame
#パラグラフの追加
pg = text_frame.add_paragraph()
#フォントサイズを大きくする
pg.font.size = Pt(24)
pg.text = "font.size = 24 pt"

#以下、フォントをMeiryo UIにしてboldにする
pg = text_frame.add_paragraph()
pg.font.bold = True
pg.font.name = "font.name = Meiryo UI"
pg.text = "font.name = Meiryo UI"

#以下、フォントをイタリックにするにする
pg = text_frame.add_paragraph()
pg.font.italic = True
pg.text = "font.italic = True"

#以下、アンダーラインを引く
pg = text_frame.add_paragraph()
pg.font.underline = True
pg.text = "font.underline = True"



shape = shapes.add_textbox(Cm(1), Cm(8), Cm(5), Cm(5))
text_frame = shape.text_frame
text_frame.clear()
pg =text_frame.paragraphs[0]
pg.text ="level　= 0"

pg =text_frame.add_paragraph()
pg.text ="level = 1"
pg.level = 1

pg =text_frame.add_paragraph()
pg.text ="level = 2"
pg.level = 2


#中央揃え
shape = shapes.add_textbox(Cm(0), Cm(11), Cm(6.14), Cm(1))
text_frame = shape.text_frame
text_frame.clear()
pg =text_frame.paragraphs[0]
pg.text ="alignment = center"
pg.alignment = PP_ALIGN.CENTER

#左揃え
shape = shapes.add_textbox(Cm(7.5), Cm(11), Cm(6.14), Cm(1))
text_frame = shape.text_frame
text_frame.clear()
pg =text_frame.paragraphs[0]
pg.text ="alignment = left "
pg.alignment = PP_ALIGN.LEFT

#右揃え
shape = shapes.add_textbox(Cm(15), Cm(11), Cm(6.14), Cm(1))
text_frame = shape.text_frame
text_frame.clear()
pg =text_frame.paragraphs[0]
pg.text ="alignment = right"
pg.alignment = PP_ALIGN.RIGHT

##########################################################################################################

# https://hogelog.com/python/python-powerpoint-5.html


slide = prs.slides.add_slide(prs.slide_layouts[6]) 

shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "テーブルの作成"
run.font.color.rgb = RGBColor(0, 0, 255)


x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
row = 4
col = 4
table = slide.shapes.add_table(row, col, x, y, cx, cy).table


#テキストの挿入
table.cell(0, 1).text = "1st Column"

#text_frameを使用して、_Paragraph.textでテキストを挿入
pg = table.cell(0,2).text_frame.paragraphs[0]
pg.text = "2nd Column"
pg.font.color.rgb = RGBColor(255, 0, 0)

#text_frameを使用して、_Run.textでテキストを挿入
pg = table.cell(0, 3).text_frame.paragraphs[0]
run = pg.add_run()
run.text = "3rd Column"
run.font.color.rgb = RGBColor(0, 255, 0)

#各セルにテキストを挿入
table.cell(1, 0).text = "1st Row"
table.cell(2, 0).text = "2nd Row"
table.cell(3, 0).text = "3rd Row"

for row in range(1, 4):
    for col in range(1, 4):
        table.cell(row, col).text = "cell({0}, {1})".format(row, col)


#最初の行を見出しとするか
table.first_col = False
#最初の列を見出しとするか
table.first_row = True

#セルを結合する
table.cell(2, 1).merge(table.cell(2,2))

#行の高さと列の幅を変える
table.rows[2].height = Cm(4)
table.columns[1].width = Cm(2)

#セルを塗りつぶす
table.cell(1, 2).fill.solid()
table.cell(1, 2).fill.fore_color.rgb = RGBColor(255, 255, 0)

#セルの余白と文字の配置を調整する
#セルの左の余白
table.cell(2, 3).margin_left = Inches(0.5)
#セルの上の余白
table.cell(2, 3).margin_top = Cm(2)
#文字の中央揃え
table.cell(1, 3).vertical_anchor = MSO_ANCHOR.MIDDLE



##########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


slide = prs.slides.add_slide(prs.slide_layouts[6]) 

shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "いろいろ"
run.font.color.rgb = RGBColor(0, 0, 255)


# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#slideshapes-objects
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype

a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(4), height=Cm(5) ,left=Cm(3), top=Cm(1))
l = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE,       width=Cm(5), height=Cm(5) ,left=Cm(5), top=Cm(5))

# https://github.com/scanny/python-pptx/blob/master/docs/dev/analysis/shp-freeform.rst

freeform_builder = slide.shapes.build_freeform(Inches(8), Inches(8))
freeform_builder.add_line_segments((
     (Inches(2),   Inches(2)),
     (Inches(1),   Inches(2)),
     (Inches(1.5), Inches(1)),
 ))
freeform_shape = freeform_builder.convert_to_shape()

##########################################################################################################
##########################################################################################################
##########################################################################################################



prs.save('connector_test3.pptx')


subprocess.Popen(['start', 'connector_test3.pptx'], shell=True)



