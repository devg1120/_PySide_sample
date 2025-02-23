
import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches


import subprocess

prs = pptx.Presentation()


#########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


slide = prs.slides.add_slide(prs.slide_layouts[6]) 


c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(4), top=Cm(2))
d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(12), top=Cm(5))

pg = c.text_frame.paragraphs[0]
pg.text ="C"
pg.alignment = PP_ALIGN.CENTER

pg = d.text_frame.paragraphs[0]
pg.text ="D"
pg.alignment = PP_ALIGN.CENTER

line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(c, 3)
line1.end_connect(d, 2)
#line1.element.spPr.prstGeom.rewrite_guides([('adj1', -635)])

e = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(4), top=Cm(10))
f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Cm(2), height=Cm(2) ,left=Cm(12), top=Cm(13))

pg = e.text_frame.paragraphs[0]
pg.text ="E"
pg.alignment = PP_ALIGN.CENTER

pg = f.text_frame.paragraphs[0]
pg.text ="F"
pg.alignment = PP_ALIGN.CENTER


line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2) )
#line1.element.flipH = True
#line1.element.flipV = True
#line1.element.rot = 16200000
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
line1.begin_connect(e, 0)
line1.end_connect(f, 0)

#https://github.com/scanny/python-pptx/blob/master/src/pptx/oxml/shapes/shared.py
#line1.element.flipH = True
#line1.element.flipV = True
#line1.element.rot = 16200000


#e.element.spPr.prstGeom.rewrite_guides([("adj1", 25000), ("adj2", 25000), ("adj3", 25000), ("adj4", 43750)])
#print(e.element.spPr.prstGeom.gd_lst)
#e.element.flipH = True
#e.element.rot = 1620
print(type(line1.element))
print(type(line1.element.nvCxnSpPr))
print(type(line1.element.nvCxnSpPr.cNvPr)) #<class 'pptx.oxml.shapes.shared.CT_NonVisualDrawingProps'>

print(type(line1.element.nvCxnSpPr.cNvCxnSpPr))  #<class 'pptx.oxml.shapes.connector.CT_NonVisualConnectorProperties'>
print(type(line1.element.nvCxnSpPr.cNvCxnSpPr.stCxn))  #<class 'pptx.oxml.shapes.connector.CT_Connection'>
print(type(line1.element.nvCxnSpPr.cNvCxnSpPr.endCxn)) #<class 'pptx.oxml.shapes.connector.CT_Connection'>

print(type(line1.element.nvCxnSpPr.nvPr))
print(type(line1.element.spPr))
print(type(line1.element.spPr.prstGeom))

#line1.element.spPr.prstGeom.rewrite_guides([('adj1', 10000)])
#line1.element.spPr.prstGeom.prst = "bentConnector4"
#line1.element.spPr.prstGeom.prst = MSO_SHAPE.MATH_EQUAL
#line1.element.spPr.prstGeom.prst = MSO_CONNECTOR_TYPE.ELBOW
#line1.element.spPr.prstGeom.prst = MSO_CONNECTOR_TYPE.MIXED

#line1.element.spPr.prstGeom.rewrite_guides([('adj1', 42857), ('adj2', 118325)])
""""
          <a:prstGeom prst="bentConnector4">
            <a:avLst>
              <a:gd name="adj1" fmla="val 42857"/>
              <a:gd name="adj2" fmla="val 118325"/>
            </a:avLst>
          </a:prstGeom>
"""

#line1.element.spPr.prstGeom.rewrite_guides([("adj1", 50000), ("adj2", 50000)])
#line1.element.spPr.prstGeom.rewrite_guides([('adj2', 3500)])

##########################################################################################################
##########################################################################################################
##########################################################################################################


prs.save('tmp/test1/test1.pptx')
prs.save('tmp/test2/test2.pptx')


subprocess.Popen(['start', 'tmp/test1/test1.pptx'], shell=True)



