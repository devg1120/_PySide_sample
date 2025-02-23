
import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches


import subprocess

prs = pptx.Presentation()


#########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


sld0 = prs.slides.add_slide(prs.slide_layouts[6]) 

# [A]図形の追加と形状の変形
#-----------------------------------------------------------------------------------------------
rect0 = sld0.shapes.add_shape(		# shapeオブジェクト➀を追加
    MSO_SHAPE.ROUNDED_RECTANGLE,   	# 図形の種類を[丸角四角形]に指定
    Cm(2), Cm(2),                  	# 挿入位置の指定　左上の座標の指定
    Cm(5), Cm(3))                  	# 挿入図形の幅と高さの指定

for i in range(3):

    rect1 = sld0.shapes.add_shape(     	# shapeオブジェクト➁を追加
        MSO_SHAPE.DECAGON,              # 図形の種類を[六角形]に指定
        Cm(2+i*5), Cm(7),
        Cm(3), Cm(3))
    
    rect1.adjustments[0] = 0.5*(i+1)   # 図形の形状を変形(高さ方向)


for i in range(3):
    rect2 = sld0.shapes.add_shape(     	# shapeオブジェクト➂を追加
        MSO_SHAPE.STAR_10_POINT,        # 図形の種類を[星(10)]に指定
        Cm(2+i*5), Cm(12),
        Cm(3), Cm(3))
    
    rect2.adjustments[0] = 0.21*(i+1)   # 図形の形状を変形
    

# [B]図形の塗り潰しとテキストの設定
#-----------------------------------------------------------------------------------------------
rect0.fill.solid()                                   # shapeオブジェクト➀を単色で塗り潰す
rect0.fill.fore_color.rgb = RGBColor(250, 100, 100)  # RGB指定で色を指定

pg = rect0.text_frame.paragraphs[0]   	# shapeオブジェクト➀のTextFrameの取得
pg.text = 'ROUNDED_RECTANGLE'         	# TextFrameにテキストを設定
pg.font.size = Pt(10)                	# テキストの文字サイズを10ポイントとする



sld1 = prs.slides.add_slide(prs.slide_layouts[6])

left = top = width = height = Inches(1)
txBox = sld1.shapes.add_textbox(left, top, width, height)    #Text Box Shapeオブジェクトの追加

#----------------------------------------------------------------------------------------------------
tf = txBox.text_frame		# TextFrameオブジェクトの設定
tf.text = "This is text inside a textbox"            # TextFrameオブジェクトにはデフォルトで1つ段落を持つ

p = tf.add_paragraph()		                           # paragraphオブジェクトの追加作成(2段落目)
p.text = "This is a second paragraph that's bold"    # textプロパティによる文字列の設定
p.font.bold = True		                               # font.boldプロパティによる太文字設定

p = tf.add_paragraph()		                      # paragraphオブジェクトの追加作成(3段落目)
p.text = "This is a third paragraph that's big"	# textプロパティによる文字列の設定
p.font.size = Pt(40)		                        # font.sizeプロパティによる文字サイズの設定



sld2 = prs.slides.add_slide(prs.slide_layouts[6]) # 空白のスライドを追加

# <画像の貼り付け> ----------------------------------------------------------------------

# 画像-➀                                            　X座標、Y座標、横幅、縦幅
pic0 = sld2.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(1), Cm(7), Cm(5))
# 画像-➁
pic1 = sld2.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(7), Cm(7), Cm(5))
# 画像-➂
pic2 = sld2.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(13), Cm(7), Cm(5))
# 画像-➃
pic3 = sld2.shapes.add_picture('sample_picture.jpg', Cm(12), Cm(4), Cm(7), Cm(5))
# 画像-➄
pic4 = sld2.shapes.add_picture('sample_picture.jpg', Cm(12), Cm(12), Cm(7), Cm(5))


# <トリミング、回転> --------------------------------------------------------------------

pic0.crop_top = 0.25     # 上から0.25(25％伸長)
pic2.crop_bottom = 0.25  # 下から0.25(25％伸長)

pic3.rotation = 45       # 45度回転(時計回り)
pic4.rotation = -45      # -45度回転

prs.save('autoshape.pptx')


subprocess.Popen(['start', 'autoshape.pptx'], shell=True)



