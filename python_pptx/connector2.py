
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Inches

prs = Presentation()

slide_layout = prs.slide_layouts[6]

slide = prs.slides.add_slide(slide_layout)

line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(2))
line1.line.fill.background()
line1.line.fill.solid()
line1.line.fill.fore_color.rgb = RGBColor(128, 255, 0)

prs.save('test.pptx')


