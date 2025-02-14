from pptx import Presentation
from pptx.util import Inches, Pt   # 単位指定用クラス(Inches:インチ単位, Pt:ポイント単位)のインポート

prs = Presentation()
sld0 = prs.slides.add_slide(prs.slide_layouts[6])

left = top = width = height = Inches(1)
txBox = sld0.shapes.add_textbox(left, top, width, height)    #Text Box Shapeオブジェクトの追加

#----------------------------------------------------------------------------------------------------
tf = txBox.text_frame		# TextFrameオブジェクトの設定
tf.text = "This is text inside a textbox"            # TextFrameオブジェクトにはデフォルトで1つ段落を持つ

p = tf.add_paragraph()		                           # paragraphオブジェクトの追加作成(2段落目)
p.text = "This is a second paragraph that's bold"    # textプロパティによる文字列の設定
p.font.bold = True		                               # font.boldプロパティによる太文字設定

p = tf.add_paragraph()		                      # paragraphオブジェクトの追加作成(3段落目)
p.text = "This is a third paragraph that's big"	# textプロパティによる文字列の設定
p.font.size = Pt(40)		                        # font.sizeプロパティによる文字サイズの設定

prs.save('Blog_テキストボックスの作成.pptx')
