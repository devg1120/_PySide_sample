import os
import signal
import psutil

import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.util import Cm,Pt,Inches

from shape_class import Shape
from shape_class import Point
from shape_class import Shape
from shape_class import Rectangle
from shape_class import RoundedRectangle
from shape_class import Cube
from shape_class import Arc
from shape_class import Oval
from shape_class import TextBox
from shape_class import Line
from shape_class import FreeForm
from shape_class import Table

#import sys
#sys.path.append("./pptx_canvas/canvas")

#from canvas  import Canvas

import subprocess

fill_color = RGBColor(100, 149, 237)
line_color = RGBColor( 65, 105, 225)
line_weight = 1
text_color = RGBColor(0,0,0)
text_size = 16
text_align =  PP_ALIGN.CENTER
text_bold = True
text_italic = False

prs = pptx.Presentation()


##########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


slide = prs.slides.add_slide(prs.slide_layouts[6]) 

shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "いろいろ"
run.font.color.rgb = RGBColor(0, 0, 255)

# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#slideshapes-objects
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype


s1 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(18), top=Cm(2), height=pptx.util.Inches(0.6))
s2 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(21), top=Cm(5), height=pptx.util.Inches(0.6))

line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line.begin_connect(s1, 3)
line.end_connect(s2, 0)

l = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE,       width=Cm(3), height=Cm(2) ,left=Cm(5), top=Cm(1))

a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Pt(80), height=Pt(40) ,left=Pt(100), top=Pt(200))

b = Rectangle(slide, 80, 40, 300,300, text = "OK")

r = Rectangle(slide, 150, 100, 10,10, text = "ok2", 
        fill_color = RGBColor.from_string("5a2bf7"), 
        line_color = RGBColor.from_string("ff0000"), 
        line_weight = 1)

o = RoundedRectangle(slide, 150, 100, 300,20, text = "RONDED", 
        fill_color = RGBColor.from_string("5a2bf7"), 
        line_color = RGBColor.from_string("ff0000"), 
        line_weight = 1, rounded = 0.3)

c = Cube(slide, 150, 100, 440,50, text = "cube", depth = 0.5)
#a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 45, radius1 = 0)
#a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 90, radius1 = 0)
#a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 135, radius1 = 0)
a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 145, radius1 = 0)
#a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 162, radius1 = 0)
#a = Arc(slide, 150, 150, 550,150, text = "arc" , radius0 = 200, radius1 = 0)

o = Oval(slide, 150, 150, 550,350, text = "oval" , fill_color =RGBColor(255,255,0))


s = Shape(slide, MSO_SHAPE.DONUT,
           100, 200, 400, 300,
           RGBColor(255,0,127), line_color, line_weight ,
           "DONUT", text_color, text_size, text_align , text_bold, text_italic)


########################################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

t = TextBox(slide, 120, 80, 400,200, "TEXTbox")

start = Point(0,0)
end   = Point(400,200)
l1 = Line(slide, start, end)

pl = [Point(300,300), 
      Point(320,420), 
      Point(350,350), 
      Point(400,500)]
l2 = Line(slide, pointlist = pl)

# https://github.com/scanny/python-pptx/blob/master/docs/dev/analysis/shp-freeform.rst

freeform_builder = slide.shapes.build_freeform(Inches(2), Inches(2))
freeform_builder.add_line_segments((
     (Inches(2),   Inches(2)),
     (Inches(1),   Inches(2)),
     (Inches(1.5), Inches(1)),
 ), close = False)
freeform_shape = freeform_builder.convert_to_shape()
freeform_shape.text = "TEXT1"

freeform_shape.line.fill.background()
freeform_shape.line.fill.solid()
freeform_shape.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
freeform_shape.fill.background()
freeform_shape.fill.solid()
freeform_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

freeform_builder = slide.shapes.build_freeform(Inches(7), Inches(7))
freeform_builder.add_line_segments((
     (Inches(7),   Inches(7)),
     (Inches(6),   Inches(7)),
     (Inches(7),   Inches(6)),
     (Inches(6.5), Inches(6)),
 ), close = False)
freeform_shape = freeform_builder.convert_to_shape()
freeform_shape.text = "TEXT2"

freeform_shape.line.fill.background()
freeform_shape.line.fill.solid()
freeform_shape.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
freeform_shape.fill.background()
freeform_shape.fill.solid()
freeform_shape.fill.fore_color.rgb = RGBColor(100,44,100)



freeform_builder = slide.shapes.build_freeform(Cm(5), Cm(10))
f = freeform_builder.add_line_segments([],close = False)
#f.add_line_segments([(Cm(5),   Cm(10))], close = False)
f.add_line_segments([(Cm(10),   Cm(10))], close = False)
f.add_line_segments([(Cm(7),   Cm(14))], close = False)
#f.add_line_segments([(Cm(9),   Cm(9))])
s = f.convert_to_shape()
s.text = "OKOKOKK"
s.line.fill.background()
s.line.fill.solid()
s.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
s.fill.background()
s.fill.solid()
s.fill.fore_color.rgb = RGBColor(100,44,100)




shp = slide.shapes.add_textbox(Pt(200), Pt(200), Pt(170), Pt(50))
shp.text = "あいうえお"

shp.text_frame.paragraphs[0].font.size = Pt(30)
shp.line.color.rgb = RGBColor(255, 0, 0)

shp.fill.solid()
shp.fill.fore_color.rgb = RGBColor(250, 250, 0)


begin_x = Pt(90)
begin_y = Pt(90)
end_x   = Inches(4.0)
end_y   = Inches(4.0)

###
# 引数
#  引数1:MSO_CONNECTOR.STRAIGHT
#  引数2: begin x 座標
#  引数3: begin y 座標
#  引数4: end   x 座標
#  引数5  end   y 座標
###
shape = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
  )

shape.line.fill.background()
shape.line.fill.solid()
shape.line.fill.fore_color.rgb = RGBColor(235,27,27)
shape.line.width = Pt(15)



#r = open("pid.txt","r")
#pid = r.read()
#r.close()
#print("kill:",pid)
#os.kill(int(pid), signal.SIGTERM)
#os.kill(int(pid), signal.CTRL_BREAK_EVENT)
#os.kill(int(pid), signal.SIGKILL)
#psutil.Process(int(pid)).terminate ()
#psutil.Process(int(19524)).terminate ()

#####################################################################################
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

table = Table(slide, 4,4,10,10, 400,300)

#テキストの挿入
table.cell(0, 1).text = "1st Column"

#text_frameを使用して、_Paragraph.textでテキストを挿入
pg = table.cell(0,2).text_frame.paragraphs[0]
pg.text = "2nd Column"
pg.font.color.rgb = RGBColor(255, 0, 0)

#text_frameを使用して、_Run.textでテキストを挿入
pg = table.cell(0, 3).text_frame.paragraphs[0]
run = pg.add_run()
run.text = "3rd Column"
run.font.color.rgb = RGBColor(0, 255, 0)

#各セルにテキストを挿入
table.cell(1, 0).text = "1st Row"
table.cell(2, 0).text = "2nd Row"
table.cell(3, 0).text = "3rd Row"

for row in range(1, 4):
    for col in range(1, 4):
        table.cell(row, col).text = "cell({0}, {1})".format(row, col)



#最初の行を見出しとするか
table.first_col(False)
#最初の列を見出しとするか
table.first_row(True)

#セルを結合する
table.cell(2, 1).merge(table.cell(2,2))

#行の高さと列の幅を変える
table.rows[2].height = Cm(4)
table.columns[1].width = Cm(2)

#セルを塗りつぶす
table.cell(1, 2).fill.solid()
table.cell(1, 2).fill.fore_color.rgb = RGBColor(255, 255, 0)

#セルの余白と文字の配置を調整する
#セルの左の余白
table.cell(2, 3).margin_left = Inches(0.5)
#セルの上の余白
table.cell(2, 3).margin_top = Cm(2)
#文字の中央揃え
table.cell(1, 3).vertical_anchor = MSO_ANCHOR.MIDDLE

#####################################################################################
"""
slide = prs.slides.add_slide(prs.slide_layouts[6]) 

canvas = Canvas(slide, prs.slide_height, prs.slide_width)
ctx = canvas.getContext('2d')

ctx.fillStyle = '#FF0000'
ctx.fillRect(Pt(10), Pt(10), Pt(150), Pt(75))
"""
#####################################################################################
prs.save('lines.pptx')

p = subprocess.Popen(['start', 'lines.pptx'], shell=True)
print("run:",p.pid)
f = open("pid.txt","w")
f.write(str(p.pid))
f.close()

