# coding: UTF-8

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# センチメートルをPowerPoint上の距離に変換する関数
def Centis(length):
    centi = Inches(length/2.54)
    return centi


# ---------------------------------------------------
"テンプレート、出力ファイル名の設定"
# ---------------------------------------------------
# templateとなるpptxファイルを指定する。
template_path = "./template.pptx"
# 出力するpptxファイルを指定する。(存在しない場合、自動作成されます)
save_path = "./output.pptx"


# ---------------------------------------------------
"空のスライドの挿入"
# ---------------------------------------------------
presentation = Presentation(template_path)
title_slide_layout = presentation.slide_layouts[1]  # レイアウトや書式を元ファイルから参照する
slide = presentation.slides.add_slide(title_slide_layout)
shapes = slide.shapes


# ---------------------------------------------------
"タイトルテキストの挿入"
# ---------------------------------------------------
# 入力したい文字列
slide_title = "TEST TITLE"
shapes.title.text = slide_title


# ---------------------------------------------------
"imageの挿入"
# ---------------------------------------------------
# 挿入する位置
pic_left = Centis(1)
pic_top = Centis(5)

# imageの高さを指定
pic_height = Centis(7.9)

image_path = "./sample_graph.png"
slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)


# ---------------------------------------------------
"tableの挿入"
# ---------------------------------------------------
# 入力したいtable状のデータ
sample_table = [ ["1.1","1.2","1.3"]
                ,["2.1","2.2","2.3"]
                ,["3.1","3.2","3.3"]]

# cell内のフォントサイズ
cell_font = 20

# 挿入する位置
table_left = Centis(9.4)
table_top = Centis(5)

# tableの幅と高さ（仮）
table_width = Centis(15)
table_height = Centis(10)

# tableの行数と列数(tableのサイズ)
rows = len(sample_table)
cols = len(sample_table[0])

table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

# 表の各セルの中身を記入
for i in range(rows):
    for j in range(cols):
        cell = table.cell(i, j)
        cell.text = sample_table[i][j]
        cell.text_frame.paragraphs[0].font.size = Pt(cell_font)

# tableの高さを再調整
table.rows[0].height = Centis(1.5)
table.rows[1].height = Centis(4.9)
table.rows[2].height = Centis(1.5)

# tableの幅を再調整
table.columns[0].width = Centis((15) / 3)
table.columns[1].width = Centis((15) / 3)
table.columns[2].width = Centis((15) / 3)


# ---------------------------------------------------
"テキストボックスの挿入"
# ---------------------------------------------------
# 文字列
sample_str = "Test Text"

# テキストボックスの位置
text_left = Centis(1)
text_top = Centis(13.4)

# テキストボックスの幅と高さ
text_width = Centis(25.4-2)
text_height = Centis(5)

# 文字のフォントサイズ
text_font = 20

# 塗りつぶし色指定(R, G, B)
color = RGBColor(150, 255, 255)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

text_box.text = sample_str
text_box.text_frame.add_paragraph().font.size = Pt(text_font)
text_box.fill.solid()
text_box.fill.fore_color.rgb = color


# ---------------------------------------------------
"ファイル保存"
# ---------------------------------------------------
presentation.save(save_path)
