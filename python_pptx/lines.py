import os
import signal
import psutil

import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.util import Cm,Pt,Inches


import subprocess

prs = pptx.Presentation()


##########################################################################################################
# https://qiita.com/Mt_SQ/items/9025d26b5709ca0648c6


slide = prs.slides.add_slide(prs.slide_layouts[6]) 

shapes = slide.shapes
shape = shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(5))
text_frame = shape.text_frame

pg = text_frame.paragraphs[0]
run = pg.add_run()
run.text = "いろいろ"
run.font.color.rgb = RGBColor(0, 0, 255)

class Point():
   def  __init__(self,  x = 0 , y = 0 ):
       self.x = x
       self.y = y

class Rectangle():
   def  __init__(self, slide = None , width = 0 , height = 0 , left = 0, top = 0):
         self.slide  = slide
         self.width  = width
         self.height = height
         self.left   = left
         self.top    = top

         if not slide == None:
            self.shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                 width=Pt(width), height=Pt(height) ,left=Pt(left), top=Pt(top))
            self.shape.fill.solid()
            self.shape.fill.fore_color.rgb = RGBColor(250, 100, 100)

class TextBox():
   def  __init__(self, slide = None , width = 0 , height = 0 , left = 0, top = 0, text = "NoSet"):
         self.slide  = slide
         self.width  = width
         self.height = height
         self.left   = left
         self.top    = top
         self.text = text

         if not slide == None:
            self.shape = slide.shapes.add_textbox( 
                 width=Pt(width), height=Pt(height) ,left=Pt(left), top=Pt(top))
            self.shape.fill.solid()
            self.shape.fill.fore_color.rgb = RGBColor(250, 100, 100)
            self.shape.text = self.text

class Line():
   def  __init__(self, slide = None , start = None, end = None, pointlist = None):
         self.slide  = slide
         self.start  = start
         self.end    = end
         self.shape = None
         self.shapes = []

         if slide == None :
             retutn
         if pointlist == None:
            self.shape = slide.shapes.add_connector(
                     #MSO_CONNECTOR.STRAIGHT, 
                     #MSO_CONNECTOR.ELBOW, 
                     MSO_CONNECTOR.CURVE, 
                    Pt(start.x), 
                    Pt(start.y),
                    Pt(end.x ), 
                    Pt(end.y ),
                 )
            self.shape.line.fill.background()
            self.shape.line.fill.solid()
            self.shape.line.fill.fore_color.rgb = RGBColor(128, 255, 0)
            self.shape.line.width = Pt(5)
            self.shape.text = "LINE1"
         else:
             for i, point in enumerate(pointlist[:-1]):
                 next_point = pointlist[i + 1]
                 shape = slide.shapes.add_connector(
                     MSO_CONNECTOR.STRAIGHT, 
                     Pt(point.x), #begin_x, 
                     Pt(point.y), #begin_y, 
                     Pt(next_point.x), #end_x, 
                     Pt(next_point.y)  #end_y
                   )
                 self.shapes.append(shape)
             self.group = slide.shapes.add_group_shape(self.shapes)


class FreeForm():
    def __int__(self, slide, pointlist = [], text = ""):
        if len(pointlist) < 2:
            return
        x = pointlist[0].x
        y = pointlist[0].y
        freeform_builder = slide.shapes.build_freeform(Pt(x),Pt(y))
        for i, point in enumerate(pointlist):
             freeform_builder.add_line_segments((
                     Pt(point.x), 
                     Pt(point.y)  
                   ), clise = False)
        self.freeform_shape = freeform_builder.convert_to_shape()
        self.freeform_shape.text = "free TEXT1"


# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#slideshapes-objects
# https://python-pptx.readthedocs.io/en/latest/api/enum/MsoAutoShapeType.html#msoautoshapetype


s1 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(18), top=Cm(2), height=pptx.util.Inches(0.6))
s2 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(21), top=Cm(5), height=pptx.util.Inches(0.6))

line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line.begin_connect(s1, 3)
line.end_connect(s2, 0)

l = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE,       width=Cm(3), height=Cm(2) ,left=Cm(5), top=Cm(1))

a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,       width=Pt(80), height=Pt(40) ,left=Pt(100), top=Pt(200))

b = Rectangle(slide, 80, 40, 300,300)
t = TextBox(slide, 120, 80, 400,300, "TEXTbox")

start = Point(0,0)
end   = Point(400,200)
l1 = Line(slide, start, end)

pl = [Point(300,300), 
      Point(320,420), 
      Point(350,350), 
      Point(400,500)]
l2 = Line(slide, pointlist = pl)

# https://github.com/scanny/python-pptx/blob/master/docs/dev/analysis/shp-freeform.rst

freeform_builder = slide.shapes.build_freeform(Inches(2), Inches(2))
freeform_builder.add_line_segments((
     (Inches(2),   Inches(2)),
     (Inches(1),   Inches(2)),
     (Inches(1.5), Inches(1)),
 ), close = False)
freeform_shape = freeform_builder.convert_to_shape()
freeform_shape.text = "TEXT1"

freeform_shape.line.fill.background()
freeform_shape.line.fill.solid()
freeform_shape.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
freeform_shape.fill.background()
freeform_shape.fill.solid()
freeform_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

freeform_builder = slide.shapes.build_freeform(Inches(7), Inches(7))
freeform_builder.add_line_segments((
     (Inches(7),   Inches(7)),
     (Inches(6),   Inches(7)),
     (Inches(7),   Inches(6)),
     (Inches(6.5), Inches(6)),
 ), close = False)
freeform_shape = freeform_builder.convert_to_shape()
freeform_shape.text = "TEXT2"

freeform_shape.line.fill.background()
freeform_shape.line.fill.solid()
freeform_shape.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
freeform_shape.fill.background()
freeform_shape.fill.solid()
freeform_shape.fill.fore_color.rgb = RGBColor(100,44,100)

##########################################################################################################



"""
freeform_builder.add_line_segments(
 vertices = [
   (Cm(10),   Cm(10)),
   (Cm(9),   Cm(10)),
   (Cm(9.5), Cm(9)),
 ],
 close = False
)
freeform_builder.convert_to_shape()
"""

freeform_builder = slide.shapes.build_freeform(Cm(5), Cm(10))
f = freeform_builder.add_line_segments([],close = False)
#f.add_line_segments([(Cm(5),   Cm(10))], close = False)
f.add_line_segments([(Cm(10),   Cm(10))], close = False)
f.add_line_segments([(Cm(7),   Cm(14))], close = False)
#f.add_line_segments([(Cm(9),   Cm(9))])
s = f.convert_to_shape()
s.text = "OKOKOKK"
s.line.fill.background()
s.line.fill.solid()
s.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
s.fill.background()
s.fill.solid()
s.fill.fore_color.rgb = RGBColor(100,44,100)




"""
freeform_builder = shapes.build_freeform(
    1500, 1000, scale=1000/Inches(1)
)
freeform_builder.add_line_segments((
    (2000, 2000),
    (1000, 2000),
    (1500, 1000),
))
freeform_shape = freeform_builder.convert_to_shape()
"""
##########################################################################################################

shp = slide.shapes.add_textbox(Pt(200), Pt(200), Pt(170), Pt(50))
shp.text = "あいうえお"

shp.text_frame.paragraphs[0].font.size = Pt(30)
shp.line.color.rgb = RGBColor(255, 0, 0)

shp.fill.solid()
shp.fill.fore_color.rgb = RGBColor(250, 250, 0)

##########################################################################################################

begin_x = Pt(90)
begin_y = Pt(90)
end_x   = Inches(4.0)
end_y   = Inches(4.0)

###
# 引数
#  引数1:MSO_CONNECTOR.STRAIGHT
#  引数2: begin x 座標
#  引数3: begin y 座標
#  引数4: end   x 座標
#  引数5  end   y 座標
###
shape = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
  )

shape.line.fill.background()
shape.line.fill.solid()
shape.line.fill.fore_color.rgb = RGBColor(235,27,27)
shape.line.width = Pt(15)


r = open("pid.txt","r")
pid = r.read()
r.close()
print("kill:",pid)
#os.kill(int(pid), signal.SIGTERM)
#os.kill(int(pid), signal.CTRL_BREAK_EVENT)
#os.kill(int(pid), signal.SIGKILL)
#psutil.Process(int(pid)).terminate ()
#psutil.Process(int(19524)).terminate ()


prs.save('lines.pptx')

p = subprocess.Popen(['start', 'lines.pptx'], shell=True)
print("run:",p.pid)
f = open("pid.txt","w")
f.write(str(p.pid))
f.close()

