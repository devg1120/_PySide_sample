
import pptx
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.util import Cm

prs = pptx.Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6]) # this yields a blank slide in my default template

s1 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(2), top=Cm(2), height=pptx.util.Inches(1))
s2 = slide.shapes.add_picture('./sample_picture.jpg', left=Cm(4), top=Cm(5), height=pptx.util.Inches(1))

line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW, Cm(2), Cm(2), Cm(2), Cm(2))
line.begin_connect(s1, 3)
line.end_connect(s2, 0)

prs.save('connector_test.pptx')

