from pptx import Presentation
from pptx.util import Cm         # 単位指定用クラス Cm


prs = Presentation()
sld0 = prs.slides.add_slide(prs.slide_layouts[6]) # 空白のスライドを追加

# <画像の貼り付け> ----------------------------------------------------------------------

# 画像-➀                                            　X座標、Y座標、横幅、縦幅
pic0 = sld0.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(1), Cm(7), Cm(5))
# 画像-➁
pic1 = sld0.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(7), Cm(7), Cm(5))
# 画像-➂
pic2 = sld0.shapes.add_picture('sample_picture.jpg', Cm(1), Cm(13), Cm(7), Cm(5))
# 画像-➃
pic3 = sld0.shapes.add_picture('sample_picture.jpg', Cm(12), Cm(4), Cm(7), Cm(5))
# 画像-➄
pic4 = sld0.shapes.add_picture('sample_picture.jpg', Cm(12), Cm(12), Cm(7), Cm(5))


# <トリミング、回転> --------------------------------------------------------------------

pic0.crop_top = 0.25     # 上から0.25(25％伸長)
pic2.crop_bottom = 0.25  # 下から0.25(25％伸長)

pic3.rotation = 45       # 45度回転(時計回り)
pic4.rotation = -45      # -45度回転


prs.save('Shapeオブジェクト(画像)_apply2.pptx')
