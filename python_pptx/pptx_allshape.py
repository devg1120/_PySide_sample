from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Cm, Pt

OUTPUT_PATH = "shapetest.pptx"

ppt = Presentation()
SLIDE_WIDTH = ppt.slide_width
SLIDE_HEIGHT = ppt.slide_height
plot_top = Cm(2)
plotzone_height = SLIDE_HEIGHT-plot_top


def table_x(x_posi,x_size):	# 0~
	sukima = (SLIDE_WIDTH/6 - x_size) / 2
	left = SLIDE_WIDTH/6 * x_posi + sukima
	return left
def table_y(y_posi):		# 0~
	sukima = Cm(0.2)
	top = plot_top+(plotzone_height/5)*y_posi+sukima
	return top
def text_x(x_posi):
    return SLIDE_WIDTH/6*x_posi
def text_y(y_posi):
    return plot_top+(plotzone_height/5)*(y_posi+1)-Cm(1)

# Page 1
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)

# 直線
shapes.add_shape(MSO_SHAPE.LINE_INVERSE,		table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))

# 矢印
shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW,	table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(1))
shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW,		table_x(0,Cm(1)),	table_y(2), Cm(1), Cm(2))
shapes.add_shape(MSO_SHAPE.QUAD_ARROW,			table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))

shapes.add_shape(MSO_SHAPE.LEFT_ARROW,			table_x(1, Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,			table_x(1, Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.UP_ARROW,			table_x(1, Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DOWN_ARROW,			table_x(1, Cm(2)),	table_y(3), Cm(2), Cm(2))
# BENT ARROW
shapes.add_shape(MSO_SHAPE.BENT_ARROW,			table_x(2, Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.BENT_UP_ARROW,		table_x(2, Cm(2)),	table_y(1), Cm(2), Cm(2))
# CIRCULAR_ARROW
shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW,		table_x(2, Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_CIRCULAR_ARROW,	table_x(2, Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_CIRCULAR_ARROW,	table_x(2, Cm(2)),	table_y(4), Cm(2), Cm(2))
# curved arrow
shapes.add_shape(MSO_SHAPE.CURVED_LEFT_ARROW,   table_x(3, Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CURVED_RIGHT_ARROW,  table_x(3, Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CURVED_UP_ARROW,     table_x(3, Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CURVED_DOWN_ARROW,   table_x(3, Cm(2)),	table_y(3), Cm(2), Cm(2))
# 丸(楕円)
shapes.add_shape(MSO_SHAPE.OVAL,				table_x(3, Cm(2)),	table_y(4), Cm(2), Cm(2))
# 長方形
shapes.add_shape(MSO_SHAPE.RECTANGLE,			table_x(4, Cm(2)),	table_y(0), Cm(2), Cm(2))
# CUBE
shapes.add_shape(MSO_SHAPE.CUBE,				table_x(4, Cm(2)),	table_y(1), Cm(2), Cm(2))
# 三角
shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,		table_x(4, Cm(2)),	table_y(2), Cm(2), Cm(2))
# 五角形
shapes.add_shape(MSO_SHAPE.PENTAGON,			table_x(4, Cm(2)),	table_y(3), Cm(2), Cm(2))
# 五角形
shapes.add_shape(MSO_SHAPE.HEXAGON,				table_x(4, Cm(2)),	table_y(4), Cm(2), Cm(2))


# 角が丸い四角形
shapes.add_shape(MSO_SHAPE.ROUND_1_RECTANGLE,			table_x(5, Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,		table_x(5, Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ROUND_2_SAME_RECTANGLE,		table_x(5, Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,			table_x(5, Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,	table_x(5, Cm(2)),	table_y(4), Cm(2), Cm(2))

mso_list = [
	["LINE_INVERSE","LEFT_ARROW","BENT_ARROW","CURVED_LEFT_ARROW","RECTANGLE","ROUND_1_RECTANGLE"],
	["LEFT_RIGHT_ARROW","RIGHT_ARROW","BENT_UP_ARROW","CURVED_RIGHT_ARROW","CUBE","ROUND_2_DIAG_RECTANGLE"],
	["UP_DOWN_ARROW","UP_ARROW","CIRCULAR_ARROW","CURVED_UP_ARROW","RIGHT_TRIANGLE","ROUND_2_SAME_RECTANGLE"],
	["QUAD_ARROW","DOWN_ARROW","LEFT_CIRCULAR_ARROW","CURVED_DOWN_ARROW","PENTAGON","ROUNDED_RECTANGLE"],
	["","","LEFT_RIGHT_CIRCULAR_ARROW","OVAL","HEXAGON","ROUNDED_RECTANGULAR\v_CALLOUT"],
]
title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List MAIN"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE		# 図形内のテキストに合わせて、必要に応じてフォントサイズを縮小


# Page 2
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)


shapes.add_shape(MSO_SHAPE.ARC,					table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.BALLOON,				table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.BEVEL,				table_x(0,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.BLOCK_ARC,			table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CAN,					table_x(0,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CHART_PLUS,			table_x(1,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CHART_STAR,			table_x(1,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CHART_X,				table_x(1,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CHEVRON,				table_x(1,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CHORD,				table_x(1,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CLOUD,				table_x(2,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CLOUD_CALLOUT,		table_x(2,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CORNER,				table_x(2,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CORNER_TABS,			table_x(2,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CROSS,				table_x(2,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CURVED_DOWN_RIBBON,	table_x(3,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.CURVED_UP_RIBBON,	table_x(3,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DECAGON,				table_x(3,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DIAGONAL_STRIPE,		table_x(3,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DIAMOND,				table_x(3,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DODECAGON,			table_x(4,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DONUT,				table_x(4,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DOUBLE_BRACE,		table_x(4,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DOUBLE_BRACKET,		table_x(4,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DOUBLE_WAVE,			table_x(4,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.DOWN_RIBBON,			table_x(5,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.EXPLOSION1,			table_x(5,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.EXPLOSION2,			table_x(5,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,		table_x(5,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FRAME,				table_x(5,Cm(2)),	table_y(4), Cm(2), Cm(2))
mso_list = [
	["ARC"		,"CHART_PLUS"	,"CLOUD"		,"CURVED_DOWN_RIBBON"	,"DODECAGON"		,"DOWN_RIBBON"	],
	["BALLOON"	,"CHART_STAR"	,"CLOUD_CALLOUT","CURVED_UP_RIBBON"		,"DONUT"			,"EXPLOSION1"	],
	["BEVEL"	,"CHART_X"		,"CORNER"		,"DECAGON"				,"DOUBLE_BRACE"		,"EXPLOSION2"	],
	["BLOCK_ARC","CHERVRON"		,"CORNER_TABS"	,"DIAGONAL_STRIPE"		,"DOUBLE_BRACKET"	,"FOLDED_CORNER"],
	["CAN"		,"CHORD"		,"CROSS"		,"DIAMOND"				,"DOUBLE_WAVE"		,"FRAME"		],
]

title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List (1)"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE		# 図形内のテキストに合わせて、必要に応じてフォントサイズを縮小

# Page 3
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)


shapes.add_shape(MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS			,table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_CARD						,table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_COLLATE					,table_x(0,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_CONNECTOR					,table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DATA						,table_x(0,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DECISION					,table_x(1,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DELAY						,table_x(1,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DIRECT_ACCESS_STORAGE		,table_x(1,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DISPLAY					,table_x(1,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_DOCUMENT					,table_x(1,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_EXTRACT					,table_x(2,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE			,table_x(2,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_MAGNETIC_DISK				,table_x(2,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_MANUAL_INPUT				,table_x(2,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_MANUAL_OPERATION			,table_x(2,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_MERGE						,table_x(3,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_MULTIDOCUMENT				,table_x(3,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFLINE_STORAGE			,table_x(3,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR			,table_x(3,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_OR							,table_x(3,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS			,table_x(4,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_PREPARATION				,table_x(4,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_PROCESS					,table_x(4,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_PUNCHED_TAPE				,table_x(4,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_SEQUENTIAL_ACCESS_STORAGE	,table_x(4,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_SORT						,table_x(5,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_STORED_DATA				,table_x(5,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION			,table_x(5,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FLOWCHART_TERMINATOR					,table_x(5,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.FOLDED_CORNER						,table_x(5,Cm(2)),	table_y(4), Cm(2), Cm(2))


mso_list = [
	["FLOWCHART_\vALTERNATE_PROCESS","FLOWCHART_DECISION"				,"FLOWCHART_EXTRACT"			,"FLOWCHART_MERGE"				,"FLOWCHART\v_PREDEFINED_PROCESS"			,"FLOWCHART_SORT"				],
	["FLOWCHART_CARD"				,"FLOWCHART_DELAY"					,"FLOWCHART_INTERNAL_STORAGE"	,"FLOWCHART_MULTIDOCUMENT"		,"FLOWCHART_PREPERATION"					,"FLOWCHART_STORED_DATA"		],
	["FLOWCHART_COLLATE"			,"FLOWCHART_DIRECT_\vACCESS_STORAGE"	,"FLOWCHART_MAGNETIC_DISK"		,"FLOWCHART_OFFLINE_STRAGE"		,"FLOWCHART_PROCESS"						,"FLOWCHART\v_SUMMING_JUNCTION"	],
	["FLOWCHART_CONNECTOR"			,"FLOWCHART_DISPLAY"				,"FLOWCHART_MANUAL_INPUT"		,"FLOWCHART_\vOFFPAGE_CONNECTOR","FLOWCHART_PUNCHED_TAPE"					,"FLOWCHART_TERMINATOR"			],
	["FLOWCHART_DATA"				,"FLOWCHART_DOCUMENT"				,"FLOWCHART\v_MANUAL_OPERATION"	,"FLOWCHART_OR"					,"FLOWCHART_SEQUENTIAL\v_ACCESS_STORAGE"	,"FOLDED_CORNER"				],
]

title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List (2) FLOWCHART"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE		# 図形内のテキストに合わせて、必要に応じてフォントサイズを縮小


# Page 4
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)


shapes.add_shape(MSO_SHAPE.FUNNEL					,table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.GEAR_6					,table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.GEAR_9					,table_x(0,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.HALF_FRAME				,table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.HEART					,table_x(0,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.HEPTAGON					,table_x(1,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.HORIZONTAL_SCROLL		,table_x(1,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE		,table_x(1,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_BRACE				,table_x(1,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_BRACKET				,table_x(1,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_RIBBON		,table_x(2,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT			,table_x(2,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_DIVIDE				,table_x(2,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_EQUAL				,table_x(2,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_MINUS				,table_x(2,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY			,table_x(3,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_NOT_EQUAL			,table_x(3,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MATH_PLUS				,table_x(3,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.MOON						,table_x(3,Cm(2)),	table_y(3), Cm(2), Cm(2))
#shapes.add_shape(MSO_SHAPE.NO_SYMBOL				,table_x(3,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.NON_ISOSCELES_TRAPEZOID	,table_x(4,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.NOTCHED_RIGHT_ARROW		,table_x(4,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.OCTAGON					,table_x(4,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.PARALLELOGRAM			,table_x(4,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.PIE						,table_x(4,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.PIE_WEDGE				,table_x(5,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.PLAQUE					,table_x(5,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.PLAQUE_TABS				,table_x(5,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.REGULAR_PENTAGON			,table_x(5,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SMILEY_FACE				,table_x(5,Cm(2)),	table_y(4), Cm(2), Cm(2))


mso_list = [
	["FUNNEL"		,"HEPTAGON"				,"LEFT_RIGHT_RIBBON","MATH_MULTIPLY"	,"NON_ISOSCELES_TRAPEZOID"	,"PIE_WEDGE"		],
	["GEAR_6"		,"HORIZONTAL_SCROLL"	,"LIGHTNING_BOLT"	,"MATH_NOT_EQUAL"	,"NOTCHED_RIGHT_ARROW"		,"PLAQUE"			],
	["GEAR_9"		,"ISOSCELES_TRIANGLE"	,"MATH_DIVIDE"		,"MATH_PLUS"		,"OCTAGON"					,"PLAQUE_TABS"		],
	["HALF_FRAME"	,"LEFT_BRACE"			,"MATH_EQUAL"		,"MOON"				,"PARALLELOGRAM"			,"REGULAR_PENTAGON"	],
	["HEART"		,"LEFT_BRACKET"			,"MATH_MINUS"		,"NO_SYMBOL"		,"PIE"						,"SMILEY_FACE"		],
]

title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List (3)"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE		# 図形内のテキストに合わせて、必要に応じてフォントサイズを縮小

# Page 5
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)


shapes.add_shape(MSO_SHAPE.RIGHT_BRACE				,table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.RIGHT_BRACKET			,table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE			,table_x(0,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SNIP_1_RECTANGLE			,table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SNIP_2_DIAG_RECTANGLE	,table_x(0,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SNIP_2_SAME_RECTANGLE	,table_x(1,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SNIP_ROUND_RECTANGLE		,table_x(1,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SQUARE_TABS				,table_x(1,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_10_POINT			,table_x(1,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_12_POINT			,table_x(1,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_16_POINT			,table_x(2,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_24_POINT			,table_x(2,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_32_POINT			,table_x(2,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_4_POINT				,table_x(2,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_5_POINT				,table_x(2,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_6_POINT				,table_x(3,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_7_POINT				,table_x(3,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STAR_8_POINT				,table_x(3,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.STRIPED_RIGHT_ARROW		,table_x(3,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SUN						,table_x(3,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.SWOOSH_ARROW				,table_x(4,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.TEAR						,table_x(4,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.TRAPEZOID				,table_x(4,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.U_TURN_ARROW				,table_x(4,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.UP_RIBBON				,table_x(4,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.VERTICAL_SCROLL			,table_x(5,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.WAVE						,table_x(5,Cm(2)),	table_y(1), Cm(2), Cm(2))

mso_list = [
	["RIGHT_BRACE"			,"SNIP_2_SAME_RECTANGLE","STAR_16_POINT","STAR_6_POINT"			,"SWOOSH_ARROW"	,"VERTICAL_SCROLL"],
	["RIGHT_BRACKET"		,"SNIP_ROUND_RECTANGLE"	,"STAR_24_POINT","STAR_7_POINT"			,"TEAR"			,"WAVE"],
	["RIGHT_TRIANGLE"		,"SQUARE_TABS"			,"STAR_32_POINT","STAR_8_POINT"			,"TRAPEZOID"	,""],
	["SNIP_1_RECTANGLE"		,"STAR_10_POINT"		,"STAR_4_POINT"	,"STRIPED_RIGHT_ARROW"	,"U_TURN_ARROW"	,""],
	["SNIP_2_DIAG_RECTANGLE","STAR_12_POINT"		,"STAR_5_POINT"	,"SUN"					,"UP_RIBBON"	,""],
]

title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List (4)"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE		# 図形内のテキストに合わせて、必要に応じてフォントサイズを縮小

# Page 5
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
shapes = slide.shapes
# set line
# horizontal line
for i in range(6):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, 0, plot_top+(plotzone_height/5)*i, SLIDE_WIDTH, 0)
# virtical line
for i in range(5):
    shapes.add_shape(MSO_SHAPE.LINE_INVERSE, SLIDE_WIDTH/6*(i+1), plot_top, 0, plotzone_height)


shapes.add_shape(MSO_SHAPE.DOWN_ARROW_CALLOUT					,table_x(0,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_ARROW_CALLOUT					,table_x(0,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.UP_ARROW_CALLOUT						,table_x(0,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.RIGHT_ARROW_CALLOUT					,table_x(0,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW_CALLOUT				,table_x(0,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW_CALLOUT				,table_x(1,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1						,table_x(1,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_ACCENT_BAR			,table_x(1,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_BORDER_AND_ACCENT_BAR	,table_x(1,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER				,table_x(1,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2						,table_x(2,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR			,table_x(2,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2_BORDER_AND_ACCENT_BAR	,table_x(2,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_2_NO_BORDER				,table_x(2,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_3						,table_x(2,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_3_ACCENT_BAR			,table_x(3,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_3_BORDER_AND_ACCENT_BAR	,table_x(3,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_3_NO_BORDER				,table_x(3,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_4						,table_x(3,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_4_ACCENT_BAR			,table_x(3,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR	,table_x(4,Cm(2)),	table_y(0), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_4_NO_BORDER				,table_x(4,Cm(2)),	table_y(1), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.OVAL_CALLOUT							,table_x(4,Cm(2)),	table_y(2), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.QUAD_ARROW_CALLOUT					,table_x(4,Cm(2)),	table_y(3), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT					,table_x(4,Cm(2)),	table_y(4), Cm(2), Cm(2))
shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT			,table_x(5,Cm(2)),	table_y(0), Cm(2), Cm(2))

mso_list = [
	["DOWN_ARROW_CALLOUT"		,"UP_DOWN_ARROW_CALLOUT"					,"LINE_CALLOUT_2"						,"LINE_CALLOUT_3_ACCENT_BAR"		,"LINE_CALLOUT_4\v_BORDER_AND_ACCENT_BAR","RPINDED_RECTANGULAR\v_CALLOUT"],
	["LEFT_ARROW_CALLOUT"		,"LINE_CALLOUT_1"							,"LINE_CALLOUT_2_ACCENT_BAR"			,"LINE_CALLOUT_3\v_BORDER_AND_ACCENT_BAR","LINE_CALLOUT_4_NO_BORDER",""],
	["UP_ARROW_CALLOUT"			,"LINE_CALLOUT_1_ACCENT_BAR"				,"LINE_CALLOUT_2\v_BORDERED_AND_ACCENT_BAR","LINE_CALLOUT_3_NO_BORDER","OVAL_CALLOUT",""],
	["RIGHT_ARROW_CALLOUT"		,"LINE_CALLOUT_1\v_BORDERED_AND_ACCENT_BAR"	,"LINE_CALLOUT_2_NO_BORDER"				,"LINE_CALLOUT_4","QUAD_ARROW_CALLOUT"								,""],
	["LEFT_RIGHT_ARROW_CALLOUT"	,"LINE_CALLOUT_1_NO_BORDER"					,"LINE_CALLOUT_3" 						,"LINE_CALLOUT_4_ACCENT_BAR","RECTANGULAR_CALLOUT"					,""],
]

title_width = Cm(25)
title_left = (SLIDE_WIDTH - title_width) / 2
textbox = shapes.add_textbox(title_left, 0, title_width, plot_top)
textbox.text = "MSO_SHAPE List (5) callout"
textbox.text_frame.paragraphs[0].font.size = Pt(30)
textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
textbox.text_frame.virtical_anchor = MSO_ANCHOR.MIDDLE

for row in range(5):
	for col in range(6):
		textbox = shapes.add_textbox(text_x(col), text_y(row), SLIDE_WIDTH/6, Cm(1))
		textbox.text = mso_list[row][col]
		textbox.text_frame.paragraphs[0].font.size = Pt(9)
		textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

ppt.save(OUTPUT_PATH)
